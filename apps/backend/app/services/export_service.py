from io import BytesIO
from datetime import UTC, datetime
from typing import Any
from uuid import UUID
from xml.sax.saxutils import escape

from fastapi import HTTPException, status
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import ListFlowable, ListItem, PageBreak, Paragraph, SimpleDocTemplate, Spacer
from sqlalchemy import select
from sqlalchemy.orm import Session

from app.models.generated_content import GeneratedContent
from app.models.project import Project
from app.models.user import User
from app.services.coverage_plan_service import (
    build_course_structure_shape_from_plan,
    get_approved_plan,
    list_lesson_scripts_from_plan,
)
from app.services.project_service import get_project_by_id

FOOTER_TEXT = "Virtus et Veritas Engine"


def safe_text(value: Any, fallback: str = "Nao informado") -> str:
    if value is None:
        return fallback
    if isinstance(value, str):
        return value.strip() or fallback
    if isinstance(value, bool | int | float):
        return str(value)
    if isinstance(value, list):
        text = "\n".join(safe_text(item, "") for item in value).strip()
        return text or fallback
    if isinstance(value, dict):
        title = (
            value.get("title")
            or value.get("subtitle")
            or value.get("concept")
            or value.get("question")
            or value.get("pergunta")
            or value.get("answer")
            or value.get("resposta")
        )
        description = (
            value.get("explanation")
            or value.get("feedback")
            or value.get("comment")
            or value.get("description")
            or value.get("content")
            or value.get("text")
            or value.get("summary")
            or value.get("details")
            or value.get("speaker_notes")
            or value.get("visual_suggestion")
            or value.get("interaction_question")
        )
        if title and description:
            return f"{safe_text(title, '')}: {safe_text(description, '')}".strip(": ")
        if description:
            return safe_text(description, fallback)
        if title:
            return safe_text(title, fallback)

        parts = [safe_text(item, "") for item in value.values()]
        text = "\n".join(part for part in parts if part).strip()
        return text or fallback

    return str(value)


def as_list(value: Any) -> list[Any]:
    if isinstance(value, list):
        return value
    if value is None or value == "":
        return []
    return [value]


def get_latest_presentation_deck(db: Session, project: Project) -> GeneratedContent:
    content = db.execute(
        select(GeneratedContent)
        .where(
            GeneratedContent.project_id == project.id,
            GeneratedContent.organization_id == project.organization_id,
            GeneratedContent.content_type == "presentation_deck",
        )
        .order_by(GeneratedContent.version.desc(), GeneratedContent.created_at.desc())
    ).scalars().first()

    if content is None or not content.content_json:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="A apresentacao ainda nao foi gerada para este projeto.",
        )

    return content


def get_content_number(content: GeneratedContent, field: str) -> int:
    content_json = content.content_json or {}
    script = content_json.get("lesson_script", {}) if isinstance(content_json.get("lesson_script"), dict) else {}
    quiz = content_json.get("module_quiz", {}) if isinstance(content_json.get("module_quiz"), dict) else {}
    metadata = content_json.get("metadata", {}) if isinstance(content_json.get("metadata"), dict) else {}
    for source in (script, quiz, metadata):
        try:
            return int(source.get(field) or 0)
        except (TypeError, ValueError):
            continue
    return 0


def get_lesson_scripts(db: Session, project: Project) -> list[GeneratedContent]:
    coverage_plan = get_approved_plan(db, project.id)
    contents = list_lesson_scripts_from_plan(db, coverage_plan) if coverage_plan is not None else []

    if not contents:
        contents = list(
            db.execute(
                select(GeneratedContent)
                .where(
                    GeneratedContent.project_id == project.id,
                    GeneratedContent.organization_id == project.organization_id,
                    GeneratedContent.content_type.in_(["lesson_script", "lesson_scripts"]),
                )
                .order_by(GeneratedContent.created_at.asc())
            ).scalars().all()
        )

    contents.sort(
        key=lambda content: (
            get_content_number(content, "module_number") or 9999,
            get_content_number(content, "lesson_number") or 9999,
            content.created_at,
        )
    )

    if not contents:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Nenhum roteiro de aula foi encontrado para este projeto.",
        )

    return contents


def get_module_quizzes(db: Session, project: Project) -> list[GeneratedContent]:
    contents = list(
        db.execute(
            select(GeneratedContent)
            .where(
                GeneratedContent.project_id == project.id,
                GeneratedContent.organization_id == project.organization_id,
                GeneratedContent.content_type.in_(["module_quiz", "module_quizzes"]),
            )
            .order_by(GeneratedContent.created_at.asc())
        ).scalars().all()
    )

    contents.sort(key=lambda content: (get_content_number(content, "module_number") or 9999, content.created_at))

    if not contents:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Nenhum quiz foi encontrado para este projeto.",
        )

    return contents


def get_complementary_materials(db: Session, project: Project) -> list[GeneratedContent]:
    contents = list(
        db.execute(
            select(GeneratedContent)
            .where(
                GeneratedContent.project_id == project.id,
                GeneratedContent.organization_id == project.organization_id,
                GeneratedContent.content_type.in_(["complementary_material", "complementary_materials"]),
            )
            .order_by(GeneratedContent.created_at.asc())
        ).scalars().all()
    )

    contents.sort(
        key=lambda content: (
            safe_text(
                (content.content_json or {}).get("complementary_material", {}).get("material_title")
                if isinstance((content.content_json or {}).get("complementary_material"), dict)
                else content.title,
                "",
            ).lower(),
            content.created_at,
        )
    )

    if not contents:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Nenhum material complementar foi encontrado para este projeto.",
        )

    return contents


def get_optional_latest_content(db: Session, project: Project, content_type: str) -> GeneratedContent | None:
    return db.execute(
        select(GeneratedContent)
        .where(
            GeneratedContent.project_id == project.id,
            GeneratedContent.organization_id == project.organization_id,
            GeneratedContent.content_type == content_type,
        )
        .order_by(GeneratedContent.version.desc(), GeneratedContent.created_at.desc())
    ).scalars().first()


def get_optional_contents(db: Session, project: Project, content_types: list[str]) -> list[GeneratedContent]:
    return list(
        db.execute(
            select(GeneratedContent)
            .where(
                GeneratedContent.project_id == project.id,
                GeneratedContent.organization_id == project.organization_id,
                GeneratedContent.content_type.in_(content_types),
            )
            .order_by(GeneratedContent.created_at.asc())
        ).scalars().all()
    )


def paragraph(text: Any, style: ParagraphStyle) -> Paragraph:
    return Paragraph(escape(safe_text(text)), style)


def labeled_paragraph(label: str, value: Any, label_style: ParagraphStyle, body_style: ParagraphStyle) -> list[Any]:
    return [
        Paragraph(escape(label), label_style),
        paragraph(value, body_style),
        Spacer(1, 0.18 * cm),
    ]


def footer(canvas: Any, doc: SimpleDocTemplate) -> None:
    canvas.saveState()
    canvas.setFont("Helvetica", 8)
    canvas.setFillColor(colors.HexColor("#6B7280"))
    canvas.drawString(doc.leftMargin, 1.05 * cm, FOOTER_TEXT)
    canvas.drawRightString(A4[0] - doc.rightMargin, 1.05 * cm, str(doc.page))
    canvas.restoreState()


def build_presentation_pdf(project: Project, presentation: GeneratedContent) -> bytes:
    deck = presentation.content_json or {}
    slides = as_list(deck.get("slides"))

    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=1.6 * cm,
        leftMargin=1.6 * cm,
        topMargin=1.6 * cm,
        bottomMargin=1.6 * cm,
        title=safe_text(deck.get("presentation_title"), project.title),
    )

    base_styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "VveTitle",
        parent=base_styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=24,
        leading=30,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#111827"),
        spaceAfter=16,
    )
    subtitle_style = ParagraphStyle(
        "VveSubtitle",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=12,
        leading=17,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#4B5563"),
        spaceAfter=10,
    )
    section_style = ParagraphStyle(
        "VveSection",
        parent=base_styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=15,
        leading=20,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=10,
        spaceAfter=8,
    )
    slide_title_style = ParagraphStyle(
        "VveSlideTitle",
        parent=base_styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=18,
        leading=23,
        textColor=colors.HexColor("#111827"),
        spaceAfter=6,
    )
    label_style = ParagraphStyle(
        "VveLabel",
        parent=base_styles["BodyText"],
        fontName="Helvetica-Bold",
        fontSize=9,
        leading=12,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=7,
        spaceAfter=2,
    )
    body_style = ParagraphStyle(
        "VveBody",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#1F2937"),
        spaceAfter=4,
    )
    bullet_style = ParagraphStyle(
        "VveBullet",
        parent=body_style,
        leftIndent=12,
        firstLineIndent=0,
        spaceAfter=3,
    )

    story: list[Any] = [
        Spacer(1, 2.4 * cm),
        paragraph(project.title, subtitle_style),
        paragraph(deck.get("presentation_title") or project.title, title_style),
        paragraph("Apresentacao gerada pelo Virtus et Veritas Engine", subtitle_style),
        Spacer(1, 0.5 * cm),
    ]
    story.extend(labeled_paragraph("Publico-alvo", deck.get("target_audience"), label_style, body_style))
    story.extend(labeled_paragraph("Duracao estimada", deck.get("estimated_duration"), label_style, body_style))
    story.extend(labeled_paragraph("Estilo visual sugerido", deck.get("visual_style"), label_style, body_style))
    story.extend(labeled_paragraph("Objetivo da apresentacao", deck.get("presentation_objective"), label_style, body_style))

    if slides:
        story.append(PageBreak())

    for index, slide in enumerate(slides):
        slide_data = slide if isinstance(slide, dict) else {"title": slide}
        slide_number = safe_text(slide_data.get("slide_number") or index + 1)
        story.append(paragraph(f"Slide {slide_number}", section_style))
        story.append(paragraph(slide_data.get("title"), slide_title_style))
        if slide_data.get("subtitle"):
            story.append(paragraph(slide_data.get("subtitle"), body_style))

        bullets = as_list(slide_data.get("bullets"))
        if bullets:
            story.append(Paragraph("Pontos principais", label_style))
            story.append(
                ListFlowable(
                    [ListItem(paragraph(item, bullet_style), leftIndent=10) for item in bullets[:5]],
                    bulletType="bullet",
                    start="circle",
                    leftIndent=16,
                )
            )
            story.append(Spacer(1, 0.15 * cm))

        story.extend(labeled_paragraph("Notas do apresentador", slide_data.get("speaker_notes"), label_style, body_style))
        story.extend(labeled_paragraph("Sugestao visual", slide_data.get("visual_suggestion"), label_style, body_style))
        if slide_data.get("interaction_question"):
            story.extend(
                labeled_paragraph(
                    "Pergunta de interacao",
                    slide_data.get("interaction_question"),
                    label_style,
                    body_style,
                )
            )

        if index < len(slides) - 1:
            story.append(PageBreak())

    if deck.get("closing_message"):
        story.append(PageBreak())
        story.append(paragraph("Encerramento", section_style))
        story.append(paragraph(deck.get("closing_message"), body_style))

    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    return buffer.getvalue()


def add_pptx_textbox(
    slide: Any,
    text: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    bold: bool = False,
    color: str = "1F2937",
    align: Any = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph_item = frame.paragraphs[0]
    if align is not None:
        paragraph_item.alignment = align
    run = paragraph_item.add_run()
    run.text = safe_text(text, "")
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_pptx_bullets(
    slide: Any,
    bullets: list[Any],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 20,
    max_visible: int = 6,
) -> list[Any]:
    visible = bullets[:max_visible]
    overflow = bullets[max_visible:]
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(visible):
        paragraph_item = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph_item.text = safe_text(bullet, "")
        paragraph_item.level = 0
        paragraph_item.font.size = Pt(font_size)
        paragraph_item.font.color.rgb = RGBColor.from_string("1F2937")
    return overflow


def add_pptx_footer(slide: Any) -> None:
    add_pptx_textbox(slide, FOOTER_TEXT, 0.55, 7.0, 5.5, 0.25, font_size=9, color="6B7280")


def add_pptx_note_box(slide: Any, title: str, value: Any, top: float) -> None:
    if value in (None, "", []):
        return
    add_pptx_textbox(slide, title, 0.65, top, 2.4, 0.25, font_size=9, bold=True, color="8A6A16")
    add_pptx_textbox(slide, value, 0.65, top + 0.25, 12.0, 0.55, font_size=10, color="4B5563")


def build_presentation_pptx(project: Project, presentation: GeneratedContent) -> bytes:
    deck = presentation.content_json or {}
    slides = as_list(deck.get("slides"))
    if not slides:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="A apresentacao ainda nao possui slides para exportar.",
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide_background(slide: Any) -> None:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string("F8FAFC")

    cover = prs.slides.add_slide(blank_layout)
    add_slide_background(cover)
    add_pptx_textbox(
        cover,
        deck.get("presentation_title") or project.title,
        0.85,
        1.55,
        11.6,
        1.0,
        font_size=34,
        bold=True,
        color="111827",
        align=PP_ALIGN.CENTER,
    )
    add_pptx_textbox(cover, project.title, 1.7, 2.75, 10.0, 0.45, font_size=18, color="8A6A16", align=PP_ALIGN.CENTER)
    if deck.get("estimated_duration"):
        add_pptx_textbox(
            cover,
            f"Duracao estimada: {safe_text(deck.get('estimated_duration'), '')}",
            2.4,
            3.45,
            8.5,
            0.35,
            font_size=14,
            color="4B5563",
            align=PP_ALIGN.CENTER,
        )
    add_pptx_footer(cover)

    context = prs.slides.add_slide(blank_layout)
    add_slide_background(context)
    add_pptx_textbox(context, "Objetivo e contexto", 0.65, 0.55, 11.8, 0.6, font_size=28, bold=True, color="111827")
    add_pptx_note_box(context, "Publico-alvo", deck.get("target_audience"), 1.55)
    add_pptx_note_box(context, "Objetivo da apresentacao", deck.get("presentation_objective"), 2.65)
    add_pptx_note_box(context, "Estilo visual sugerido", deck.get("visual_style"), 4.15)
    add_pptx_footer(context)

    for index, raw_slide in enumerate(slides, start=1):
        slide_data = raw_slide if isinstance(raw_slide, dict) else {"title": raw_slide}
        slide = prs.slides.add_slide(blank_layout)
        add_slide_background(slide)
        title = slide_data.get("title") or f"Slide {slide_data.get('slide_number') or index}"
        add_pptx_textbox(slide, title, 0.65, 0.45, 12.0, 0.55, font_size=28, bold=True, color="111827")
        if slide_data.get("subtitle"):
            add_pptx_textbox(slide, slide_data.get("subtitle"), 0.7, 1.08, 11.8, 0.4, font_size=15, color="4B5563")
        overflow = add_pptx_bullets(slide, as_list(slide_data.get("bullets")), 0.85, 1.65, 7.3, 3.6, font_size=20)
        if slide_data.get("visual_suggestion"):
            add_pptx_textbox(slide, "Sugestao visual", 8.45, 1.65, 3.8, 0.25, font_size=10, bold=True, color="8A6A16")
            add_pptx_textbox(slide, slide_data.get("visual_suggestion"), 8.45, 1.95, 3.8, 1.25, font_size=12, color="4B5563")
        if slide_data.get("interaction_question"):
            add_pptx_textbox(slide, "Pergunta de interacao", 8.45, 3.55, 3.8, 0.25, font_size=10, bold=True, color="8A6A16")
            add_pptx_textbox(slide, slide_data.get("interaction_question"), 8.45, 3.85, 3.8, 0.95, font_size=12, color="1F2937")
        notes = []
        if slide_data.get("speaker_notes"):
            notes.append(f"Notas do apresentador: {safe_text(slide_data.get('speaker_notes'), '')}")
        if overflow:
            notes.append("Bullets adicionais: " + "; ".join(safe_text(item, "") for item in overflow))
        if notes:
            add_pptx_textbox(slide, "\n".join(notes), 0.8, 5.85, 11.75, 0.75, font_size=9, color="6B7280")
        add_pptx_footer(slide)

    closing = prs.slides.add_slide(blank_layout)
    add_slide_background(closing)
    add_pptx_textbox(closing, "Encerramento", 0.8, 1.35, 11.7, 0.65, font_size=30, bold=True, color="111827", align=PP_ALIGN.CENTER)
    add_pptx_textbox(
        closing,
        deck.get("closing_message") or "Material gerado e organizado pelo Virtus et Veritas Engine.",
        1.5,
        2.45,
        10.3,
        1.2,
        font_size=20,
        color="1F2937",
        align=PP_ALIGN.CENTER,
    )
    add_pptx_textbox(closing, "Material gerado e organizado pelo Virtus et Veritas Engine.", 1.9, 4.1, 9.5, 0.45, font_size=13, color="6B7280", align=PP_ALIGN.CENTER)
    add_pptx_footer(closing)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def build_lesson_scripts_pdf(project: Project, lesson_scripts: list[GeneratedContent]) -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=1.7 * cm,
        leftMargin=1.7 * cm,
        topMargin=1.6 * cm,
        bottomMargin=1.6 * cm,
        title=f"Roteiros de Aula - {project.title}",
    )

    base_styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "VveLessonTitle",
        parent=base_styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=24,
        leading=30,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#111827"),
        spaceAfter=16,
    )
    subtitle_style = ParagraphStyle(
        "VveLessonSubtitle",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=12,
        leading=17,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#4B5563"),
        spaceAfter=10,
    )
    lesson_title_style = ParagraphStyle(
        "VveLessonHeading",
        parent=base_styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=18,
        leading=23,
        textColor=colors.HexColor("#111827"),
        spaceAfter=8,
    )
    section_style = ParagraphStyle(
        "VveLessonSection",
        parent=base_styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=13,
        leading=18,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=9,
        spaceAfter=5,
    )
    label_style = ParagraphStyle(
        "VveLessonLabel",
        parent=base_styles["BodyText"],
        fontName="Helvetica-Bold",
        fontSize=9,
        leading=12,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=6,
        spaceAfter=2,
    )
    body_style = ParagraphStyle(
        "VveLessonBody",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=15,
        textColor=colors.HexColor("#1F2937"),
        spaceAfter=4,
    )

    generated_at = datetime.now(UTC).strftime("%d/%m/%Y")
    story: list[Any] = [
        Spacer(1, 2.4 * cm),
        paragraph(project.title, subtitle_style),
        paragraph("Roteiros de Aula", title_style),
        paragraph(f"Gerado em {generated_at}", subtitle_style),
        paragraph("Documento exportado pelo Virtus et Veritas Engine", subtitle_style),
        PageBreak(),
    ]

    section_keys = [
        ("Objetivo da aula", ["learning_objective", "lesson_objective", "objective"]),
        ("Abertura", ["opening", "hook"]),
        ("Introducao", ["introduction", "intro"]),
        ("Desenvolvimento", ["development", "lesson_development"]),
        ("Exemplo pratico", ["practical_example", "examples", "example"]),
        ("Atividade pratica", ["practical_activity", "activity", "exercise"]),
        ("Pergunta de reflexao", ["reflection_question"]),
        ("Conclusao", ["conclusion", "closing"]),
        ("Call to action", ["call_to_action", "cta"]),
        ("Notas do instrutor", ["instructor_notes", "teaching_notes"]),
        ("Texto de narracao", ["narration", "voiceover"]),
        ("Sugestao visual", ["visual_suggestion"]),
    ]

    for index, content in enumerate(lesson_scripts):
        content_json = content.content_json or {}
        script = content_json.get("lesson_script", {}) if isinstance(content_json.get("lesson_script"), dict) else {}
        module_number = safe_text(script.get("module_number") or get_content_number(content, "module_number"), "")
        lesson_number = safe_text(script.get("lesson_number") or get_content_number(content, "lesson_number"), "")
        lesson_title = script.get("lesson_title") or content.title or f"Aula {lesson_number or index + 1}"

        story.append(paragraph(f"Modulo {module_number or '-'} | Aula {lesson_number or index + 1}", section_style))
        story.append(paragraph(lesson_title, lesson_title_style))
        story.extend(labeled_paragraph("Titulo do curso", script.get("course_title") or project.title, label_style, body_style))
        story.extend(labeled_paragraph("Modulo", script.get("module_title"), label_style, body_style))

        for label, keys in section_keys:
            value = next((script.get(key) for key in keys if script.get(key) not in (None, "", [])), None)
            if value not in (None, "", []):
                story.extend(labeled_paragraph(label, value, label_style, body_style))

        blocks = (
            as_list(script.get("main_script"))
            or as_list(script.get("sections"))
            or as_list(script.get("blocks"))
            or as_list(script.get("content_blocks"))
        )
        if blocks:
            story.append(Paragraph("Blocos da aula", section_style))
            for block_index, block in enumerate(blocks, start=1):
                block_data = block if isinstance(block, dict) else {"content": block}
                block_title = (
                    block_data.get("section_title")
                    or block_data.get("title")
                    or block_data.get("name")
                    or f"Bloco {block_index}"
                )
                story.append(Paragraph(escape(safe_text(block_title)), label_style))
                block_fields = [
                    block_data.get("narration"),
                    block_data.get("content"),
                    block_data.get("text"),
                    block_data.get("description"),
                ]
                block_text = next((value for value in block_fields if value not in (None, "", [])), "")
                story.append(paragraph(block_text, body_style))
                if block_data.get("teaching_notes"):
                    story.extend(labeled_paragraph("Notas de ensino", block_data.get("teaching_notes"), label_style, body_style))
                if block_data.get("visual_suggestion"):
                    story.extend(labeled_paragraph("Sugestao visual", block_data.get("visual_suggestion"), label_style, body_style))
                story.append(Spacer(1, 0.12 * cm))

        if index < len(lesson_scripts) - 1:
            story.append(PageBreak())

    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    return buffer.getvalue()


def first_quiz_value(source: dict[str, Any], keys: list[str]) -> Any:
    for key in keys:
        value = source.get(key)
        if value not in (None, "", []):
            return value
    return None


def build_quizzes_pdf(project: Project, quizzes: list[GeneratedContent]) -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=1.7 * cm,
        leftMargin=1.7 * cm,
        topMargin=1.6 * cm,
        bottomMargin=1.6 * cm,
        title=f"Quizzes do Curso - {project.title}",
    )

    base_styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "VveQuizTitle",
        parent=base_styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=24,
        leading=30,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#111827"),
        spaceAfter=16,
    )
    subtitle_style = ParagraphStyle(
        "VveQuizSubtitle",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=12,
        leading=17,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#4B5563"),
        spaceAfter=10,
    )
    quiz_title_style = ParagraphStyle(
        "VveQuizHeading",
        parent=base_styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=18,
        leading=23,
        textColor=colors.HexColor("#111827"),
        spaceAfter=8,
    )
    section_style = ParagraphStyle(
        "VveQuizSection",
        parent=base_styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=13,
        leading=18,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=9,
        spaceAfter=5,
    )
    label_style = ParagraphStyle(
        "VveQuizLabel",
        parent=base_styles["BodyText"],
        fontName="Helvetica-Bold",
        fontSize=9,
        leading=12,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=6,
        spaceAfter=2,
    )
    body_style = ParagraphStyle(
        "VveQuizBody",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=15,
        textColor=colors.HexColor("#1F2937"),
        spaceAfter=4,
    )
    option_style = ParagraphStyle(
        "VveQuizOption",
        parent=body_style,
        leftIndent=12,
        firstLineIndent=0,
        spaceAfter=3,
    )

    generated_at = datetime.now(UTC).strftime("%d/%m/%Y")
    story: list[Any] = [
        Spacer(1, 2.4 * cm),
        paragraph(project.title, subtitle_style),
        paragraph("Quizzes do Curso", title_style),
        paragraph(f"Gerado em {generated_at}", subtitle_style),
        paragraph("Documento exportado pelo Virtus et Veritas Engine", subtitle_style),
        PageBreak(),
    ]

    for quiz_index, content in enumerate(quizzes):
        content_json = content.content_json or {}
        quiz = content_json.get("module_quiz", {}) if isinstance(content_json.get("module_quiz"), dict) else {}
        module_number = safe_text(quiz.get("module_number") or get_content_number(content, "module_number"), "")
        module_title = quiz.get("module_title") or content.title or "Quiz"
        quiz_title = quiz.get("quiz_title") or quiz.get("title") or f"Quiz - Modulo {module_number or quiz_index + 1}"

        story.append(paragraph(quiz_title, quiz_title_style))
        story.extend(labeled_paragraph("Modulo", module_number or "Nao informado", label_style, body_style))
        story.extend(labeled_paragraph("Titulo do modulo", module_title, label_style, body_style))
        if quiz.get("instructions"):
            story.extend(labeled_paragraph("Instrucoes", quiz.get("instructions"), label_style, body_style))

        questions = (
            as_list(quiz.get("questions"))
            or as_list(quiz.get("perguntas"))
            or as_list(quiz.get("items"))
        )
        if questions:
            story.append(Paragraph("Perguntas", section_style))

        for question_index, question in enumerate(questions, start=1):
            question_data = question if isinstance(question, dict) else {"question": question}
            question_text = first_quiz_value(question_data, ["question", "pergunta", "title", "text", "content"])
            story.append(Paragraph(escape(f"Pergunta {question_index}"), label_style))
            story.append(paragraph(question_text, body_style))

            options = (
                as_list(question_data.get("options"))
                or as_list(question_data.get("alternatives"))
                or as_list(question_data.get("alternativas"))
                or as_list(question_data.get("answers"))
            )
            if options:
                story.append(Paragraph("Alternativas", label_style))
                story.append(
                    ListFlowable(
                        [
                            ListItem(
                                paragraph(
                                    f"{safe_text(option.get('letter') or chr(65 + option_index), '')}. "
                                    f"{safe_text(option.get('text') or option.get('content') or option.get('value'), '')}"
                                    if isinstance(option, dict)
                                    else f"{chr(65 + option_index)}. {safe_text(option, '')}",
                                    option_style,
                                ),
                                leftIndent=10,
                            )
                            for option_index, option in enumerate(options)
                        ],
                        bulletType="bullet",
                        start="circle",
                        leftIndent=16,
                    )
                )
                story.append(Spacer(1, 0.12 * cm))

            correct_answer = first_quiz_value(
                question_data,
                ["correct_answer", "answer", "resposta_correta", "correct_option", "resposta"],
            )
            explanation = first_quiz_value(
                question_data,
                ["explanation", "feedback", "comment", "comentario", "explicacao"],
            )
            if correct_answer not in (None, "", []):
                story.extend(labeled_paragraph("Resposta correta", correct_answer, label_style, body_style))
            if explanation not in (None, "", []):
                story.extend(labeled_paragraph("Explicacao", explanation, label_style, body_style))
            story.append(Spacer(1, 0.2 * cm))

        if quiz_index < len(quizzes) - 1:
            story.append(PageBreak())

    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    return buffer.getvalue()


def build_complementary_materials_pdf(project: Project, materials: list[GeneratedContent]) -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=1.7 * cm,
        leftMargin=1.7 * cm,
        topMargin=1.6 * cm,
        bottomMargin=1.6 * cm,
        title=f"Materiais Complementares - {project.title}",
    )

    base_styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "VveMaterialTitle",
        parent=base_styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=24,
        leading=30,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#111827"),
        spaceAfter=16,
    )
    subtitle_style = ParagraphStyle(
        "VveMaterialSubtitle",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=12,
        leading=17,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#4B5563"),
        spaceAfter=10,
    )
    material_title_style = ParagraphStyle(
        "VveMaterialHeading",
        parent=base_styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=18,
        leading=23,
        textColor=colors.HexColor("#111827"),
        spaceAfter=8,
    )
    section_style = ParagraphStyle(
        "VveMaterialSection",
        parent=base_styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=13,
        leading=18,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=9,
        spaceAfter=5,
    )
    label_style = ParagraphStyle(
        "VveMaterialLabel",
        parent=base_styles["BodyText"],
        fontName="Helvetica-Bold",
        fontSize=9,
        leading=12,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=6,
        spaceAfter=2,
    )
    body_style = ParagraphStyle(
        "VveMaterialBody",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=15,
        textColor=colors.HexColor("#1F2937"),
        spaceAfter=4,
    )

    generated_at = datetime.now(UTC).strftime("%d/%m/%Y")
    story: list[Any] = [
        Spacer(1, 2.4 * cm),
        paragraph(project.title, subtitle_style),
        paragraph("Materiais Complementares", title_style),
        paragraph(f"Gerado em {generated_at}", subtitle_style),
        paragraph("Documento exportado pelo Virtus et Veritas Engine", subtitle_style),
        PageBreak(),
    ]

    list_sections = [
        ("Aplicacoes praticas", ["practical_applications", "applications", "aplicacoes"]),
        ("Exercicios reflexivos", ["reflection_exercises", "exercises", "perguntas_reflexivas"]),
        ("Proximos passos", ["recommended_next_steps", "next_steps", "proximos_passos"]),
    ]
    extra_fields = [
        ("Resumo", ["summary"]),
        ("Detalhes", ["details"]),
        ("Conteudo", ["content", "text", "description"]),
    ]

    for index, content in enumerate(materials):
        content_json = content.content_json or {}
        material = (
            content_json.get("complementary_material", {})
            if isinstance(content_json.get("complementary_material"), dict)
            else {}
        )
        material_title = material.get("material_title") or material.get("title") or content.title or f"Material {index + 1}"
        material_type = material.get("material_type") or material.get("type") or "Material complementar"

        story.append(paragraph(material_title, material_title_style))
        story.extend(labeled_paragraph("Tipo do material", material_type, label_style, body_style))
        if material.get("overview"):
            story.extend(labeled_paragraph("Visao geral", material.get("overview"), label_style, body_style))

        concepts = as_list(material.get("key_concepts") or material.get("concepts") or material.get("conceitos_chave"))
        if concepts:
            story.append(Paragraph("Conceitos-chave", section_style))
            for concept_index, concept in enumerate(concepts, start=1):
                concept_data = concept if isinstance(concept, dict) else {"concept": concept}
                concept_title = (
                    concept_data.get("concept")
                    or concept_data.get("title")
                    or concept_data.get("name")
                    or f"Conceito {concept_index}"
                )
                concept_explanation = (
                    concept_data.get("explanation")
                    or concept_data.get("description")
                    or concept_data.get("text")
                    or concept_data.get("content")
                )
                story.append(Paragraph(escape(safe_text(concept_title)), label_style))
                if concept_explanation not in (None, "", []):
                    story.append(paragraph(concept_explanation, body_style))
                story.append(Spacer(1, 0.12 * cm))

        for label, keys in list_sections:
            value = next((material.get(key) for key in keys if material.get(key) not in (None, "", [])), None)
            items = as_list(value)
            if items:
                story.append(Paragraph(label, section_style))
                for item in items:
                    story.append(paragraph(item, body_style))
                story.append(Spacer(1, 0.12 * cm))

        for label, keys in extra_fields:
            value = next((material.get(key) for key in keys if material.get(key) not in (None, "", [])), None)
            if value not in (None, "", []):
                story.extend(labeled_paragraph(label, value, label_style, body_style))

        if index < len(materials) - 1:
            story.append(PageBreak())

    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    return buffer.getvalue()


def build_full_course_pdf(
    project: Project,
    course_summary: GeneratedContent | None,
    course_structure: GeneratedContent | None,
    lesson_scripts: list[GeneratedContent],
    quizzes: list[GeneratedContent],
    materials: list[GeneratedContent],
    presentation: GeneratedContent | None,
) -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=1.7 * cm,
        leftMargin=1.7 * cm,
        topMargin=1.6 * cm,
        bottomMargin=1.6 * cm,
        title=f"Curso Completo - {project.title}",
    )

    base_styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "VveFullTitle",
        parent=base_styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=24,
        leading=30,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#111827"),
        spaceAfter=16,
    )
    subtitle_style = ParagraphStyle(
        "VveFullSubtitle",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=12,
        leading=17,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#4B5563"),
        spaceAfter=10,
    )
    block_title_style = ParagraphStyle(
        "VveFullBlockTitle",
        parent=base_styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=19,
        leading=24,
        textColor=colors.HexColor("#111827"),
        spaceAfter=10,
    )
    item_title_style = ParagraphStyle(
        "VveFullItemTitle",
        parent=base_styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=14,
        leading=19,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=8,
        spaceAfter=5,
    )
    label_style = ParagraphStyle(
        "VveFullLabel",
        parent=base_styles["BodyText"],
        fontName="Helvetica-Bold",
        fontSize=9,
        leading=12,
        textColor=colors.HexColor("#8A6A16"),
        spaceBefore=6,
        spaceAfter=2,
    )
    body_style = ParagraphStyle(
        "VveFullBody",
        parent=base_styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=15,
        textColor=colors.HexColor("#1F2937"),
        spaceAfter=4,
    )

    def add_section_title(title: str) -> None:
        if story and not isinstance(story[-1], PageBreak):
            story.append(PageBreak())
        story.append(paragraph(title, block_title_style))

    def add_list_section(label: str, value: Any) -> None:
        items = as_list(value)
        if not items:
            return
        story.append(Paragraph(escape(label), label_style))
        for item in items:
            story.append(paragraph(item, body_style))
        story.append(Spacer(1, 0.12 * cm))

    generated_at = datetime.now(UTC).strftime("%d/%m/%Y")
    story: list[Any] = [
        Spacer(1, 2.4 * cm),
        paragraph(project.title, subtitle_style),
        paragraph("Curso Completo", title_style),
        paragraph(f"Gerado em {generated_at}", subtitle_style),
        paragraph("Material gerado e organizado pelo Virtus et Veritas Engine", subtitle_style),
        PageBreak(),
    ]

    if course_summary and course_summary.content_json:
        add_section_title("Resumo do curso")
        summary = course_summary.content_json.get("course_summary", {})
        if not isinstance(summary, dict):
            summary = course_summary.content_json
        summary_fields = [
            ("Titulo", ["title", "course_title"]),
            ("Descricao", ["long_description", "short_description", "description"]),
            ("Promessa", ["promise"]),
            ("Publico-alvo", ["target_audience"]),
            ("Objetivo", ["objective", "course_objective", "transformation_statement"]),
            ("Carga horaria estimada", ["estimated_duration", "workload", "duration"]),
            ("Copy de vendas", ["suggested_sales_copy"]),
        ]
        for label, keys in summary_fields:
            value = next((summary.get(key) for key in keys if summary.get(key) not in (None, "", [])), None)
            if value not in (None, "", []):
                story.extend(labeled_paragraph(label, value, label_style, body_style))
        add_list_section("Principais resultados de aprendizagem", summary.get("what_student_will_learn") or summary.get("learning_outcomes"))
        add_list_section("Diferenciais do curso", summary.get("course_differentials"))

    if course_structure and course_structure.content_json:
        add_section_title("Estrutura geral do curso")
        structure = course_structure.content_json
        course = structure.get("course", {}) if isinstance(structure.get("course"), dict) else structure
        if isinstance(course, dict):
            story.extend(labeled_paragraph("Titulo", course.get("title") or project.title, label_style, body_style))
            if course.get("description"):
                story.extend(labeled_paragraph("Descricao", course.get("description"), label_style, body_style))
            modules = as_list(course.get("modules"))
        else:
            modules = []
            story.append(paragraph(course, body_style))
        for module_index, module in enumerate(modules, start=1):
            module_data = module if isinstance(module, dict) else {"title": module}
            story.append(paragraph(f"Modulo {module_data.get('module_number') or module_index}: {safe_text(module_data.get('module_title') or module_data.get('title'), '')}", item_title_style))
            if module_data.get("description"):
                story.append(paragraph(module_data.get("description"), body_style))
            lessons = as_list(module_data.get("lessons") or module_data.get("classes") or module_data.get("aulas"))
            for lesson_index, lesson in enumerate(lessons, start=1):
                lesson_data = lesson if isinstance(lesson, dict) else {"title": lesson}
                story.extend(
                    labeled_paragraph(
                        f"Aula {lesson_data.get('lesson_number') or lesson_index}",
                        lesson_data.get("lesson_title") or lesson_data.get("title") or lesson_data,
                        label_style,
                        body_style,
                    )
                )
                if lesson_data.get("objective"):
                    story.extend(labeled_paragraph("Objetivo", lesson_data.get("objective"), label_style, body_style))

    if lesson_scripts:
        add_section_title("Roteiros de aula")
        section_keys = [
            ("Objetivo da aula", ["learning_objective", "lesson_objective", "objective"]),
            ("Abertura", ["opening", "hook"]),
            ("Desenvolvimento", ["development", "lesson_development"]),
            ("Exemplo pratico", ["practical_example", "examples", "example"]),
            ("Atividade pratica", ["practical_activity", "activity", "exercise"]),
            ("Conclusao", ["conclusion", "closing"]),
            ("Call to action", ["call_to_action", "cta"]),
        ]
        for index, content in enumerate(lesson_scripts):
            script = (content.content_json or {}).get("lesson_script", {})
            script = script if isinstance(script, dict) else {}
            story.append(paragraph(f"Modulo {safe_text(script.get('module_number'), '')} | Aula {safe_text(script.get('lesson_number') or index + 1, '')}", item_title_style))
            story.extend(labeled_paragraph("Titulo", script.get("lesson_title") or content.title, label_style, body_style))
            story.extend(labeled_paragraph("Modulo", script.get("module_title"), label_style, body_style))
            for label, keys in section_keys:
                value = next((script.get(key) for key in keys if script.get(key) not in (None, "", [])), None)
                if value not in (None, "", []):
                    story.extend(labeled_paragraph(label, value, label_style, body_style))
            for block_index, block in enumerate(as_list(script.get("main_script") or script.get("sections") or script.get("blocks")), start=1):
                block_data = block if isinstance(block, dict) else {"content": block}
                story.extend(
                    labeled_paragraph(
                        block_data.get("section_title") or block_data.get("title") or f"Bloco {block_index}",
                        block_data.get("narration") or block_data.get("content") or block_data.get("text") or block,
                        label_style,
                        body_style,
                    )
                )

    if quizzes:
        add_section_title("Quizzes")
        for quiz_index, content in enumerate(quizzes):
            quiz = (content.content_json or {}).get("module_quiz", {})
            quiz = quiz if isinstance(quiz, dict) else {}
            story.append(paragraph(quiz.get("module_title") or content.title or f"Quiz {quiz_index + 1}", item_title_style))
            if quiz.get("instructions"):
                story.extend(labeled_paragraph("Instrucoes", quiz.get("instructions"), label_style, body_style))
            for question_index, question in enumerate(as_list(quiz.get("questions") or quiz.get("perguntas")), start=1):
                question_data = question if isinstance(question, dict) else {"question": question}
                story.extend(
                    labeled_paragraph(
                        f"Pergunta {question_index}",
                        first_quiz_value(question_data, ["question", "pergunta", "title", "text", "content"]),
                        label_style,
                        body_style,
                    )
                )
                add_list_section("Alternativas", question_data.get("options") or question_data.get("alternatives") or question_data.get("alternativas"))
                correct_answer = first_quiz_value(question_data, ["correct_answer", "answer", "resposta_correta", "correct_option"])
                if correct_answer not in (None, "", []):
                    story.extend(labeled_paragraph("Resposta correta", correct_answer, label_style, body_style))
                explanation = first_quiz_value(question_data, ["explanation", "feedback", "comment", "comentario"])
                if explanation not in (None, "", []):
                    story.extend(labeled_paragraph("Explicacao", explanation, label_style, body_style))

    if materials:
        add_section_title("Materiais complementares")
        for material_index, content in enumerate(materials):
            material = (content.content_json or {}).get("complementary_material", {})
            material = material if isinstance(material, dict) else {}
            story.append(paragraph(material.get("material_title") or content.title or f"Material {material_index + 1}", item_title_style))
            story.extend(labeled_paragraph("Tipo", material.get("material_type") or "Material complementar", label_style, body_style))
            if material.get("overview"):
                story.extend(labeled_paragraph("Visao geral", material.get("overview"), label_style, body_style))
            add_list_section("Conceitos-chave", material.get("key_concepts") or material.get("concepts"))
            add_list_section("Aplicacoes praticas", material.get("practical_applications"))
            add_list_section("Exercicios reflexivos", material.get("reflection_exercises"))
            add_list_section("Proximos passos", material.get("recommended_next_steps"))

    if presentation and presentation.content_json:
        add_section_title("Apresentacao pronta")
        deck = presentation.content_json
        story.extend(labeled_paragraph("Titulo", deck.get("presentation_title") or presentation.title, label_style, body_style))
        for label, key in [
            ("Publico-alvo", "target_audience"),
            ("Duracao", "estimated_duration"),
            ("Estilo visual", "visual_style"),
            ("Objetivo", "presentation_objective"),
        ]:
            if deck.get(key):
                story.extend(labeled_paragraph(label, deck.get(key), label_style, body_style))
        for slide_index, slide in enumerate(as_list(deck.get("slides")), start=1):
            slide_data = slide if isinstance(slide, dict) else {"title": slide}
            story.append(paragraph(f"Slide {slide_data.get('slide_number') or slide_index}: {safe_text(slide_data.get('title'), '')}", item_title_style))
            add_list_section("Bullets", slide_data.get("bullets"))
            if slide_data.get("speaker_notes"):
                story.extend(labeled_paragraph("Notas do apresentador", slide_data.get("speaker_notes"), label_style, body_style))
            if slide_data.get("visual_suggestion"):
                story.extend(labeled_paragraph("Sugestao visual", slide_data.get("visual_suggestion"), label_style, body_style))
            if slide_data.get("interaction_question"):
                story.extend(labeled_paragraph("Pergunta de interacao", slide_data.get("interaction_question"), label_style, body_style))
        if deck.get("closing_message"):
            story.extend(labeled_paragraph("Mensagem de encerramento", deck.get("closing_message"), label_style, body_style))

    add_section_title("Encerramento")
    story.append(paragraph("Material gerado e organizado pelo Virtus et Veritas Engine.", body_style))

    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    return buffer.getvalue()


def export_presentation_pdf(db: Session, current_user: User, project_id: UUID) -> tuple[Project, bytes]:
    project = get_project_by_id(db, current_user, project_id)
    presentation = get_latest_presentation_deck(db, project)
    return project, build_presentation_pdf(project, presentation)


def export_presentation_pptx(db: Session, current_user: User, project_id: UUID) -> tuple[Project, bytes]:
    project = get_project_by_id(db, current_user, project_id)
    presentation = get_latest_presentation_deck(db, project)
    return project, build_presentation_pptx(project, presentation)


def export_lesson_scripts_pdf(db: Session, current_user: User, project_id: UUID) -> tuple[Project, bytes]:
    project = get_project_by_id(db, current_user, project_id)
    lesson_scripts = get_lesson_scripts(db, project)
    return project, build_lesson_scripts_pdf(project, lesson_scripts)


def export_quizzes_pdf(db: Session, current_user: User, project_id: UUID) -> tuple[Project, bytes]:
    project = get_project_by_id(db, current_user, project_id)
    quizzes = get_module_quizzes(db, project)
    return project, build_quizzes_pdf(project, quizzes)


def export_complementary_materials_pdf(db: Session, current_user: User, project_id: UUID) -> tuple[Project, bytes]:
    project = get_project_by_id(db, current_user, project_id)
    materials = get_complementary_materials(db, project)
    return project, build_complementary_materials_pdf(project, materials)


def export_full_course_pdf(db: Session, current_user: User, project_id: UUID) -> tuple[Project, bytes]:
    project = get_project_by_id(db, current_user, project_id)
    lesson_scripts = get_optional_contents(db, project, ["lesson_script", "lesson_scripts"])
    lesson_scripts.sort(
        key=lambda content: (
            get_content_number(content, "module_number") or 9999,
            get_content_number(content, "lesson_number") or 9999,
            content.created_at,
        )
    )
    quizzes = get_optional_contents(db, project, ["module_quiz", "module_quizzes"])
    quizzes.sort(key=lambda content: (get_content_number(content, "module_number") or 9999, content.created_at))
    materials = get_optional_contents(db, project, ["complementary_material", "complementary_materials"])
    materials.sort(key=lambda content: content.created_at)

    course_structure = get_optional_latest_content(db, project, "course_structure")
    if course_structure is None:
        coverage_plan = get_approved_plan(db, project.id)
        if coverage_plan is not None:
            course_structure = GeneratedContent(content_json=build_course_structure_shape_from_plan(db, coverage_plan))

    pdf_bytes = build_full_course_pdf(
        project,
        get_optional_latest_content(db, project, "course_summary"),
        course_structure,
        lesson_scripts,
        quizzes,
        materials,
        get_optional_latest_content(db, project, "presentation_deck"),
    )
    return project, pdf_bytes
